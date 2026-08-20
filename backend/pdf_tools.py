"""
PDF Tools backend service.
Provides server-side conversions and security operations that cannot run
purely in the browser. Each endpoint accepts an uploaded file (multipart)
and streams back the processed result.
"""
import os
import re
import shutil
import uuid
import subprocess
import tempfile
import difflib
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse

router = APIRouter(prefix="/api/pdf", tags=["pdf-tools"])

WORK_ROOT = Path(tempfile.gettempdir()) / "pdfpro_jobs"
WORK_ROOT.mkdir(parents=True, exist_ok=True)

MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _new_job():
    d = WORK_ROOT / uuid.uuid4().hex
    d.mkdir(parents=True, exist_ok=True)
    return d


def _cleanup(path: Path):
    shutil.rmtree(path, ignore_errors=True)


async def _save_upload(upload: UploadFile, job: Path) -> Path:
    dest = job / (upload.filename or "input")
    with open(dest, "wb") as f:
        while chunk := await upload.read(1024 * 1024):
            f.write(chunk)
    await upload.seek(0)
    return dest


def _libreoffice(src: Path, target: str, outdir: Path) -> Path:
    """Convert a file to `target` (e.g. 'pdf', 'docx') via headless LibreOffice."""
    profile = outdir / "lo_profile"
    cmd = [
        "soffice", "--headless", "--norestore", "--nolockcheck",
        f"-env:UserInstallation=file://{profile}",
        "--convert-to", target, "--outdir", str(outdir), str(src),
    ]
    proc = subprocess.run(cmd, capture_output=True, timeout=180)
    if proc.returncode != 0:
        raise HTTPException(500, f"Conversion failed: {proc.stderr.decode()[:300]}")
    ext = target.split(":")[0]
    out = outdir / (src.stem + "." + ext)
    if not out.exists():
        cands = list(outdir.glob(f"{src.stem}.*"))
        cands = [c for c in cands if c.suffix.lstrip(".") == ext]
        if not cands:
            raise HTTPException(500, "Converted file not produced.")
        out = cands[0]
    return out


def _respond(path: Path, filename: str, ext: str, job: Path, bg: BackgroundTasks):
    bg.add_task(_cleanup, job)
    return FileResponse(str(path), media_type=MIME.get(ext, "application/octet-stream"), filename=filename, background=bg)


@router.get("/health")
async def health():
    tools = {}
    for t in ["soffice", "gs", "qpdf", "tesseract", "pdftoppm"]:
        tools[t] = shutil.which(t) is not None
    try:
        import ocrmypdf  # noqa
        tools["ocrmypdf"] = True
    except Exception:
        tools["ocrmypdf"] = False
    return {"ok": True, "tools": tools}


@router.post("/inspect")
async def inspect_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """Report whether a PDF uses a legacy (non-Unicode) Devanagari font so the
    client can convert its ASCII-mapped text to Unicode before displaying/editing."""
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        legacy = _is_legacy_hindi(src)
        has_text = _has_text_layer(src)
        ratio = _devanagari_ratio(src) if has_text else 0.0
        looks_english = _looks_like_english(src) if has_text else False
        fonts = sorted(_pdf_font_names(src))[:20]
        _cleanup(job)
        return JSONResponse({
            "legacy_hindi": legacy,
            "has_text": has_text,
            "devanagari_ratio": ratio,
            "looks_english": looks_english,
            "fonts": fonts,
        })
    except Exception:
        _cleanup(job)
        return JSONResponse({"legacy_hindi": False, "has_text": False, "devanagari_ratio": 0.0, "looks_english": False, "fonts": []})


# ---------- Office / HTML -> PDF ----------
@router.post("/office-to-pdf")
async def office_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = _libreoffice(src, "pdf", job)
        name = Path(file.filename).stem + ".pdf"
        return _respond(out, name, "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


@router.post("/html-to-pdf")
async def html_to_pdf(background_tasks: BackgroundTasks, url: str = Form(None), html: str = Form(None)):
    import requests
    job = _new_job()
    try:
        content = html
        if url:
            r = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            content = r.text
        if not content:
            raise HTTPException(400, "Provide a URL or HTML content.")
        src = job / "page.html"
        src.write_text(content, encoding="utf-8")
        out = _libreoffice(src, "pdf", job)
        return _respond(out, "webpage.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


def _has_text_layer(src: Path) -> bool:
    """True if the PDF contains a meaningful extractable text layer."""
    import pdfplumber
    try:
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages[:5]:
                if len((page.extract_text() or "").strip()) > 20:
                    return True
    except Exception:
        pass
    return False


def _ocr_if_scanned(src: Path, job: Path, lang: str = "eng") -> Path:
    """If the PDF is image-only (scanned), OCR it so downstream converters see text."""
    if _has_text_layer(src):
        return src
    _require_ocr_tools()
    lang = _sanitize_lang(lang)
    import ocrmypdf
    ocred = job / (src.stem + "_ocred.pdf")
    try:
        ocrmypdf.ocr(str(src), str(ocred), skip_text=True, language=lang, optimize=0, progress_bar=False)
        if ocred.exists():
            return ocred
    except Exception:
        pass
    return src


def _run_ocr(src: Path, out_pdf: Path, lang: str, sidecar: Path = None):
    """Run ocrmypdf with image-cleaning (unpaper) + deskew for better accuracy
    on legacy/scanned Hindi forms. Falls back to a plain run if cleaning fails."""
    _require_ocr_tools()
    import ocrmypdf
    kw = dict(force_ocr=True, language=_sanitize_lang(lang), optimize=0,
              progress_bar=False, deskew=True)
    if shutil.which("unpaper"):
        kw["clean"] = True
    if sidecar is not None:
        kw["sidecar"] = str(sidecar)
    try:
        ocrmypdf.ocr(str(src), str(out_pdf), **kw)
    except Exception:
        kw.pop("clean", None)
        kw.pop("deskew", None)
        ocrmypdf.ocr(str(src), str(out_pdf), **kw)


def _clean_ocr_line(line: str) -> str:
    """Tidy up OCR noise from dotted/underscore 'fill-in' leaders in forms."""
    import re
    line = re.sub(r"[\x00-\x08\x0b\x0e-\x1f\x7f-\x9f\ufffe\uffff]", "", line)
    line = re.sub(r"[.\u2026·]{4,}", " …… ", line)   # dotted leaders
    line = re.sub(r"[_=~—–-]{4,}", " —— ", line)      # underline / dash leaders
    line = re.sub(r"[<>«»]{2,}", " ", line)            # stray arrows OCR invents
    line = re.sub(r"[ \t]{2,}", " ", line)
    return line.strip()


def _looks_like_noise(line: str) -> bool:
    """Drop lines that are essentially punctuation/symbol garbage (no real words)."""
    import re
    s = re.sub(r"\s", "", line)
    if len(s) < 2:
        return True
    meaningful = sum(1 for c in s if c.isalnum() or "\u0900" <= c <= "\u097f")
    # keep only if a reasonable share of the line is letters/digits
    return meaningful < max(2, int(0.35 * len(s)))


def _force_ocr_pdf(src: Path, job: Path, lang: str) -> Path:
    """Return an OCR'd copy of the PDF (real text layer). Used for legacy-font
    Hindi PDFs where the existing text layer is ASCII-mapped gibberish."""
    out = job / (src.stem + "_forced_ocr.pdf")
    try:
        _run_ocr(src, out, lang)
    except Exception:
        return src
    return out if out.exists() else src


def _scanned_pdf_to_docx(src: Path, job: Path, out: Path, lang: str):
    """OCR a scanned/legacy PDF and build an editable Word document from the text."""
    from docx import Document
    sidecar = job / "ocr_text.txt"
    ocred = job / (src.stem + "_ocred.pdf")
    _run_ocr(src, ocred, lang, sidecar=sidecar)
    text = sidecar.read_text(encoding="utf-8", errors="ignore") if sidecar.exists() else ""
    if not text.strip():
        raise HTTPException(422, "No readable text found in this scanned PDF.")
    doc = Document()
    pages = text.split("\f")
    for pi, page in enumerate(pages):
        for line in page.splitlines():
            clean = _clean_ocr_line(line)
            if clean and not _looks_like_noise(clean):
                doc.add_paragraph(clean)
        if pi < len(pages) - 1 and page.strip():
            doc.add_page_break()
    doc.save(str(out))


def _require_ocr_tools():
    """Ensure the OCR toolchain (Tesseract + Ghostscript) is available on PATH."""
    missing = [t for t in ("tesseract", "gs") if shutil.which(t) is None]
    if missing:
        raise HTTPException(
            503,
            "OCR is temporarily unavailable on the server (missing: "
            + ", ".join(missing) + "). Please try again shortly.",
        )


def _zip_has_parts(path: Path, required: list) -> bool:
    """A valid OOXML file is a zip that contains the required internal parts
    and is not corrupt."""
    import zipfile
    try:
        if not path.exists() or path.stat().st_size < 200:
            return False
        if not zipfile.is_zipfile(str(path)):
            return False
        with zipfile.ZipFile(str(path)) as z:
            names = set(z.namelist())
            if not set(required).issubset(names):
                return False
            if z.testzip() is not None:  # first corrupt member, if any
                return False
        return True
    except Exception:
        return False


def _docx_is_valid(path: Path) -> bool:
    """Strict validity check so we never hand back a file Word refuses to open."""
    if not _zip_has_parts(path, ["[Content_Types].xml", "word/document.xml"]):
        return False
    try:
        from docx import Document
        doc = Document(str(path))
        # touch the body so a structurally-broken document raises here
        _ = doc.paragraphs
        return True
    except Exception:
        return False


def _xlsx_is_valid(path: Path) -> bool:
    """Strict validity check for Excel output."""
    if not _zip_has_parts(path, ["[Content_Types].xml", "xl/workbook.xml"]):
        return False
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True)
        _ = wb.sheetnames
        wb.close()
        return True
    except Exception:
        return False


def _normalize_office(src: Path, ext: str, job: Path) -> Path:
    """Re-save an Office file through headless LibreOffice so it opens cleanly
    in Microsoft Word / Excel. Returns the normalized file, or the original if
    LibreOffice is unavailable or the round-trip fails."""
    if shutil.which("soffice") is None:
        return src
    try:
        norm_dir = job / f"norm_{ext}"
        norm_dir.mkdir(exist_ok=True)
        produced = _libreoffice(src, ext, norm_dir)
        final = job / (src.stem + "_ms." + ext)
        shutil.copy(str(produced), str(final))
        return final
    except Exception:
        return src


def _tess_langs() -> set:
    """Languages available to the local Tesseract install."""
    try:
        proc = subprocess.run(["tesseract", "--list-langs"], capture_output=True, timeout=30)
        lines = proc.stdout.decode(errors="ignore").splitlines()[1:]
        return {ln.strip() for ln in lines if ln.strip()}
    except Exception:
        return set()


def _sanitize_lang(lang: str) -> str:
    """Keep only languages actually installed (e.g. 'hin+eng'); default to eng."""
    avail = _tess_langs()
    req = [ln for ln in (lang or "eng").split("+") if ln]
    keep = [ln for ln in req if (not avail or ln in avail)]
    return "+".join(keep) if keep else "eng"


def _docx_text_len(path: Path) -> int:
    """Total non-whitespace characters in a .docx (used to catch 'valid but empty' output)."""
    try:
        from docx import Document
        doc = Document(str(path))
        return sum(len(p.text.strip()) for p in doc.paragraphs)
    except Exception:
        return 0


def _text_pdf_to_docx(src: Path, job: Path, out: Path):
    """Build a Unicode-safe .docx directly from the PDF text layer. This preserves
    non-Latin scripts (e.g. Hindi/Devanagari) that pdf2docx sometimes drops."""
    import re
    import pdfplumber
    from docx import Document
    illegal = re.compile(r"[\x00-\x08\x0b\x0e-\x1f\x7f-\x9f\ufffe\uffff]")
    doc = Document()
    any_text = False
    with pdfplumber.open(str(src)) as pdf:
        n = len(pdf.pages)
        for pi, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            for line in text.splitlines():
                clean = illegal.sub("", line).strip()
                if clean:
                    doc.add_paragraph(clean)
                    any_text = True
            if pi < n - 1:
                doc.add_page_break()
    if not any_text:
        raise HTTPException(422, "No extractable text found in this PDF.")
    doc.save(str(out))


# Legacy (non-Unicode) Devanagari fonts store Hindi as ASCII-mapped glyph codes,
# so the extracted "text" is gibberish like: O;fäxr :i ls  ekg@o"kksZa ls ...
# We detect these and OCR the rendered pages instead (which yields real Unicode).
LEGACY_HINDI_FONTS = (
    "kruti", "krutidev", "devlys", "dev-lys", "dev lys", "chanakya", "shusha",
    "susha", "shivaji", "agra", "walkman", "shree", "richa", "aakruti",
    "sanskrit99", "amanuj", "dvb", "dvbw", "dv-ttyogesh", "yogesh", "surekh",
    "millennium", "janmang", "naidunia", "webdunia",
)


def _pdf_font_names(src: Path) -> set:
    names = set()
    try:
        import pdfplumber
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages[:5]:
                for ch in page.chars:
                    fn = ch.get("fontname")
                    if fn:
                        names.add(fn)
    except Exception:
        pass
    return names


def _is_legacy_hindi(src: Path) -> bool:
    """True if the PDF uses a known legacy (non-Unicode) Devanagari font."""
    blob = " ".join(_pdf_font_names(src)).lower()
    return any(tok in blob for tok in LEGACY_HINDI_FONTS)


def _devanagari_ratio(src: Path) -> float:
    """Fraction of non-space characters in the text layer that are real
    Unicode Devanagari (U+0900–U+097F)."""
    total = deva = 0
    try:
        import pdfplumber
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages[:5]:
                for c in (page.extract_text() or ""):
                    if c.isspace():
                        continue
                    total += 1
                    if "\u0900" <= c <= "\u097f":
                        deva += 1
    except Exception:
        return 0.0
    return (deva / total) if total else 0.0


# A small set of very common English words. Real English prose is ~30-50%% of
# these; legacy Kruti/DevLys text (ASCII-mapped Devanagari) contains virtually
# none, so this cleanly separates "genuine English" from "legacy Hindi encoded
# as ASCII" without relying on the font name.
_COMMON_EN = {
    "the", "and", "of", "to", "in", "is", "a", "for", "that", "on", "with",
    "as", "are", "be", "this", "by", "or", "at", "it", "from", "an", "was",
    "not", "which", "have", "has", "had", "will", "would", "can", "could",
    "all", "you", "your", "we", "our", "their", "they", "he", "she", "his",
    "her", "but", "if", "so", "do", "does", "been", "were", "more", "one",
    "also", "may", "such", "its", "into", "than", "when", "who", "what",
    "how", "about", "page", "document", "name", "date", "total", "no", "yes",
    "there", "these", "some", "other", "only", "over", "then", "them", "out",
}


def _looks_like_english(src: Path) -> bool:
    """Heuristic: True if the extracted text layer reads like genuine English
    prose (many common English words). Used to avoid mistaking a normal English
    PDF for legacy ASCII-mapped Hindi when the Devanagari ratio is ~0."""
    text = ""
    try:
        import pdfplumber
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages[:5]:
                text += (page.extract_text() or "") + " "
    except Exception:
        return False
    toks = re.findall(r"[a-zA-Z]{2,}", text.lower())
    if len(toks) < 3:
        return False
    hits = sum(1 for t in toks if t in _COMMON_EN)
    return (hits / len(toks)) >= 0.12


def _kruti_pdf_to_docx(src: Path, job: Path, out: Path) -> bool:
    """Build a .docx from a legacy Kruti Dev text layer, converting each line to
    proper Unicode Devanagari. Clean and noise-free (no OCR). Returns True only
    if meaningful Devanagari was produced."""
    import pdfplumber
    from docx import Document
    from krutidev import kruti_to_unicode, devanagari_count
    doc = Document()
    total_deva = 0
    with pdfplumber.open(str(src)) as pdf:
        n = len(pdf.pages)
        for pi, page in enumerate(pdf.pages):
            raw = page.extract_text() or ""
            for line in raw.splitlines():
                conv = kruti_to_unicode(line)
                total_deva += devanagari_count(conv)
                if conv.strip():
                    doc.add_paragraph(conv)
            if pi < n - 1:
                doc.add_page_break()
    if total_deva < 5:
        return False
    doc.save(str(out))
    return True


# ---------- PDF -> Word ----------
@router.post("/pdf-to-word")
async def pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...), lang: str = Form("eng")):
    from pdf2docx import Converter
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + ".docx")
        lang = _sanitize_lang(lang)
        has_text = _has_text_layer(src)
        # Legacy non-Unicode Hindi font (Kruti Dev / DevLys / Chanakya ...) OR a
        # text layer that has almost no real Devanagari while Hindi was requested:
        # the extracted text would be ASCII-mapped gibberish, so OCR the rendered
        # pages to recover proper Unicode Devanagari.
        legacy_hindi = _is_legacy_hindi(src)
        deva_ratio = _devanagari_ratio(src) if has_text else 0.0

        if legacy_hindi:
            # Legacy non-Unicode Hindi font (Kruti Dev / DevLys / Chanakya ...):
            # the text layer is ASCII-mapped, so convert it deterministically to
            # Unicode (clean, no OCR noise); fall back to OCR only if that fails.
            made = False
            try:
                made = _kruti_pdf_to_docx(src, job, out)
            except Exception:
                made = False
            if not made:
                out.unlink(missing_ok=True)
                _scanned_pdf_to_docx(src, job, out, _sanitize_lang("hin+eng"))
        elif has_text and "hin" in lang and deva_ratio < 0.15:
            # Hindi requested but the text layer has ~no Devanagari -> OCR.
            _scanned_pdf_to_docx(src, job, out, _sanitize_lang(lang))
        elif has_text:
            # 1) pdf2docx — best layout fidelity for Latin text
            try:
                cv = Converter(str(src))
                cv.convert(str(out))
                cv.close()
            except Exception:
                out.unlink(missing_ok=True)
            # 2) if pdf2docx failed OR produced an (almost) empty doc (common with
            #    Hindi/Devanagari), rebuild a Unicode-safe docx from the text layer
            if not out.exists() or not _docx_is_valid(out) or _docx_text_len(out) < 5:
                out.unlink(missing_ok=True)
                try:
                    _text_pdf_to_docx(src, job, out)
                except Exception:
                    out.unlink(missing_ok=True)
            # 3) last resort: OCR (handles legacy/non-Unicode font encodings)
            if not out.exists() or not _docx_is_valid(out) or _docx_text_len(out) < 5:
                out.unlink(missing_ok=True)
                _scanned_pdf_to_docx(src, job, out, _sanitize_lang("hin+eng"))
        else:
            _scanned_pdf_to_docx(src, job, out, lang)
        if not out.exists() or not _docx_is_valid(out):
            raise HTTPException(500, "Conversion produced an unreadable document. Please try another file.")
        # Normalize through LibreOffice so the .docx opens cleanly in MS Word.
        final = _normalize_office(out, "docx", job)
        if _docx_is_valid(final) and _docx_text_len(final) >= max(0, _docx_text_len(out) - 2):
            out = final
        return _respond(out, src.stem + ".docx", "docx", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to Word failed: {e}")


# ---------- PDF -> Excel ----------
@router.post("/pdf-to-excel")
async def pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...), lang: str = Form("eng")):
    import pdfplumber
    from openpyxl import Workbook
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        lang = _sanitize_lang(lang)
        # Legacy (non-Unicode) Hindi font -> convert the ASCII-mapped text layer to
        # Unicode Devanagari per cell/line (clean, no OCR). Otherwise OCR only if
        # the PDF is scanned (image-only).
        convert_kruti = _is_legacy_hindi(src)
        if not convert_kruti:
            src = _ocr_if_scanned(src, job, lang)
        conv = None
        if convert_kruti:
            from krutidev import kruti_to_unicode as conv
        wb = Workbook()
        wb.remove(wb.active)
        found = False
        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Page {i+1}"[:31])
                tables = page.extract_tables()
                if tables:
                    found = True
                    for tbl in tables:
                        for row in tbl:
                            ws.append([("" if c is None else (conv(c) if conv else c)) for c in row])
                        ws.append([])
                else:
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        ws.append([conv(line) if conv else line])
        if not wb.sheetnames:
            wb.create_sheet("Sheet1")
        out = job / (src.stem + ".xlsx")
        wb.save(str(out))
        if not _xlsx_is_valid(out):
            raise HTTPException(500, "Conversion produced an unreadable spreadsheet. Please try another file.")
        # Normalize through LibreOffice so the .xlsx opens cleanly in MS Excel.
        final = _normalize_office(out, "xlsx", job)
        if _xlsx_is_valid(final):
            out = final
        return _respond(out, src.stem + ".xlsx", "xlsx", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to Excel failed: {e}")


# ---------- PDF -> PowerPoint (one image per slide) ----------
@router.post("/pdf-to-ppt")
async def pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        images = convert_from_path(str(src), dpi=120)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for idx, img in enumerate(images):
            p = job / f"s{idx}.png"
            img.save(str(p), "PNG")
            slide = prs.slides.add_slide(blank)
            iw, ih = img.size
            ratio = min(prs.slide_width / iw, prs.slide_height / ih)
            w = int(iw * ratio); h = int(ih * ratio)
            left = int((prs.slide_width - w) / 2); top = int((prs.slide_height - h) / 2)
            slide.shapes.add_picture(str(p), left, top, width=w, height=h)
        out = job / (src.stem + ".pptx")
        prs.save(str(out))
        return _respond(out, src.stem + ".pptx", "pptx", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to PowerPoint failed: {e}")


# ---------- OCR ----------
@router.post("/ocr")
async def ocr_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), lang: str = Form("eng")):
    import ocrmypdf
    job = _new_job()
    try:
        _require_ocr_tools()
        lang = _sanitize_lang(lang)
        src = await _save_upload(file, job)
        out = job / (src.stem + "_ocr.pdf")
        try:
            ocrmypdf.ocr(str(src), str(out), skip_text=True, language=lang, optimize=0, progress_bar=False)
        except Exception as oe:
            raise HTTPException(500, f"OCR failed: {str(oe)[:300]}")
        if not out.exists():
            raise HTTPException(500, "OCR produced no output.")
        return _respond(out, src.stem + "_ocr.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- Protect (encrypt) ----------
@router.post("/protect")
async def protect_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), password: str = Form(...)):
    import pikepdf
    job = _new_job()
    try:
        if not password:
            raise HTTPException(400, "Password required.")
        src = await _save_upload(file, job)
        out = job / (src.stem + "_protected.pdf")
        with pikepdf.open(str(src)) as pdf:
            pdf.save(str(out), encryption=pikepdf.Encryption(owner=password, user=password, R=6))
        return _respond(out, src.stem + "_protected.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Protect failed: {e}")


# ---------- Unlock (decrypt) ----------
@router.post("/unlock")
async def unlock_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), password: str = Form("")):
    import pikepdf
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_unlocked.pdf")
        try:
            with pikepdf.open(str(src), password=password) as pdf:
                pdf.save(str(out))
        except pikepdf.PasswordError:
            raise HTTPException(400, "Wrong password for this PDF.")
        return _respond(out, src.stem + "_unlocked.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Unlock failed: {e}")


# ---------- Repair (ghostscript rewrite) ----------
@router.post("/repair")
async def repair_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_repaired.pdf")
        cmd = ["gs", "-o", str(out), "-sDEVICE=pdfwrite", "-dPDFSTOPONERROR=false", str(src)]
        proc = subprocess.run(cmd, capture_output=True, timeout=180)
        if not out.exists():
            raise HTTPException(500, f"Repair failed: {proc.stderr.decode()[:300]}")
        return _respond(out, src.stem + "_repaired.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- PDF/A ----------
@router.post("/pdfa")
async def pdf_to_pdfa(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_pdfa.pdf")
        cmd = ["gs", "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-sColorConversionStrategy=UseDeviceIndependentColor",
               "-sDEVICE=pdfwrite", "-dPDFACompatibilityPolicy=1", f"-sOutputFile={out}", str(src)]
        proc = subprocess.run(cmd, capture_output=True, timeout=180)
        if not out.exists():
            raise HTTPException(500, f"PDF/A conversion failed: {proc.stderr.decode()[:300]}")
        return _respond(out, src.stem + "_pdfa.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- Crop (trim margins by percent) ----------
@router.post("/crop")
async def crop_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), margin: float = Form(5.0)):
    from pypdf import PdfReader, PdfWriter
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        reader = PdfReader(str(src))
        writer = PdfWriter()
        m = max(0.0, min(40.0, margin)) / 100.0
        for page in reader.pages:
            box = page.mediabox
            w = float(box.width); h = float(box.height)
            page.cropbox.lower_left = (float(box.left) + w * m, float(box.bottom) + h * m)
            page.cropbox.upper_right = (float(box.right) - w * m, float(box.top) - h * m)
            writer.add_page(page)
        out = job / (src.stem + "_cropped.pdf")
        with open(out, "wb") as f:
            writer.write(f)
        return _respond(out, src.stem + "_cropped.pdf", "pdf", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Crop failed: {e}")


# ---------- Compare two PDFs (text diff) ----------
@router.post("/compare")
async def compare_pdf(background_tasks: BackgroundTasks, file1: UploadFile = File(...), file2: UploadFile = File(...)):
    import pdfplumber
    job = _new_job()
    try:
        p1 = await _save_upload(file1, job)
        p2 = await _save_upload(file2, job)

        def text_of(p):
            lines = []
            with pdfplumber.open(str(p)) as pdf:
                for page in pdf.pages:
                    lines.extend((page.extract_text() or "").splitlines())
            return lines

        a, b = text_of(p1), text_of(p2)
        sm = difflib.SequenceMatcher(None, a, b)
        rows = []
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                for line in a[i1:i2]:
                    rows.append({"type": "equal", "text": line})
            elif tag == "replace":
                for line in a[i1:i2]:
                    rows.append({"type": "removed", "text": line})
                for line in b[j1:j2]:
                    rows.append({"type": "added", "text": line})
            elif tag == "delete":
                for line in a[i1:i2]:
                    rows.append({"type": "removed", "text": line})
            elif tag == "insert":
                for line in b[j1:j2]:
                    rows.append({"type": "added", "text": line})
        ratio = round(sm.ratio() * 100, 1)
        _cleanup(job)
        return JSONResponse({"similarity": ratio, "rows": rows[:1000]})
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Compare failed: {e}")
